"""docgen MCP server: turn analysis results into deliverables.

Tools:
  * pdf_brief(title, summary, bullets)  -> one-page PDF
  * slide_deck(title, slides)           -> .pptx deck
Files are written to ./out and the path is returned to the agent.
"""
from __future__ import annotations

import os

from mcp.server.fastmcp import FastMCP

OUT_DIR = os.environ.get("ATLAS_OUT_DIR", "out")
mcp = FastMCP("docgen")


def _ensure_out() -> None:
    os.makedirs(OUT_DIR, exist_ok=True)


@mcp.tool()
def pdf_brief(title: str, summary: str, bullets: list[str]) -> str:
    """Generate a one-page PDF brief. Returns the file path."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        HRFlowable,
        ListFlowable,
        ListItem,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    _ensure_out()
    path = os.path.join(OUT_DIR, "brief.pdf")
    styles = getSampleStyleSheet()
    navy = colors.HexColor("#102A43")
    blue = colors.HexColor("#2F80ED")
    pale_blue = colors.HexColor("#EAF3FF")
    muted = colors.HexColor("#627D98")
    title_style = ParagraphStyle(
        "AtlasTitle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=29,
        textColor=navy,
        alignment=0,
        spaceAfter=5,
    )
    eyebrow_style = ParagraphStyle(
        "AtlasEyebrow",
        parent=styles["BodyText"],
        fontName="Helvetica-Bold",
        fontSize=9,
        leading=12,
        textColor=blue,
        spaceAfter=8,
    )
    summary_style = ParagraphStyle(
        "AtlasSummary",
        parent=styles["BodyText"],
        fontSize=11,
        leading=16,
        textColor=navy,
    )
    heading_style = ParagraphStyle(
        "AtlasHeading",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=16,
        textColor=navy,
        spaceAfter=10,
    )
    bullet_style = ParagraphStyle(
        "AtlasBullet",
        parent=styles["BodyText"],
        fontSize=10.5,
        leading=15,
        textColor=navy,
    )
    doc = SimpleDocTemplate(
        path,
        pagesize=letter,
        rightMargin=0.7 * inch,
        leftMargin=0.7 * inch,
        topMargin=0.65 * inch,
        bottomMargin=0.65 * inch,
        title=title,
        author="Atlas",
    )
    summary_box = Table([[Paragraph(summary, summary_style)]], colWidths=[7.1 * inch])
    summary_box.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), pale_blue),
                ("BOX", (0, 0), (-1, -1), 0.75, colors.HexColor("#B8D7FF")),
                ("LEFTPADDING", (0, 0), (-1, -1), 16),
                ("RIGHTPADDING", (0, 0), (-1, -1), 16),
                ("TOPPADDING", (0, 0), (-1, -1), 14),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 14),
            ]
        )
    )
    story = [
        Paragraph("ATLAS / MCP-GENERATED ANALYSIS", eyebrow_style),
        Paragraph(title, title_style),
        HRFlowable(width="100%", thickness=2, color=blue, spaceBefore=5, spaceAfter=20),
        summary_box,
        Spacer(1, 24),
        Paragraph("Key findings", heading_style),
        ListFlowable(
            [ListItem(Paragraph(b, bullet_style), leftIndent=12) for b in bullets],
            bulletType="bullet",
            start="circle",
            leftIndent=18,
            bulletColor=blue,
            spaceAfter=18,
        ),
        Spacer(1, 14),
        HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#D9E2EC")),
        Spacer(1, 10),
        Paragraph(
            "Generated through Atlas's SQL, Python sandbox, and document MCP servers.",
            ParagraphStyle(
                "AtlasFooter",
                parent=styles["BodyText"],
                fontSize=8.5,
                textColor=muted,
            ),
        ),
    ]
    doc.build(story)
    return os.path.abspath(path)


@mcp.tool()
def slide_deck(title: str, slides: list[dict]) -> str:
    """Generate a .pptx deck. `slides` = [{"heading": str, "bullets": [str]}]. Returns path."""
    from pptx import Presentation
    from pptx.util import Pt

    _ensure_out()
    path = os.path.join(OUT_DIR, "deck.pptx")
    prs = Presentation()
    title_layout, bullet_layout = prs.slide_layouts[0], prs.slide_layouts[1]
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = title
    for spec in slides:
        s = prs.slides.add_slide(bullet_layout)
        s.shapes.title.text = spec.get("heading", "")
        body = s.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(spec.get("bullets", [])):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
    prs.save(path)
    return os.path.abspath(path)


if __name__ == "__main__":
    mcp.run()
